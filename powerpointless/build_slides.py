import logging
from typing import List

from pptx import Presentation as presentation
from pptx.presentation import Presentation
from pptx.slide import (
    Slide,
    SlideLayout,
    SlideMaster,
    SlidePlaceholders,
    Slides,
    SlideShapes,
)

logger = logging.getLogger(__name__)


def build_slides(template: Presentation, lines: List[str]) -> Presentation:
    try:
        slide_master: SlideMaster = template.slide_master
    except Exception as e:
        raise RuntimeError(
            "Provided presentation must have at least one slide master attached"
        ) from e

    try:
        layout: SlideLayout = slide_master.slide_layouts[0]
    except Exception as e:
        raise RuntimeError("Slide master must have at least one layout") from e

    slides: Slides = template.slides

    for line in lines:
        new_slide: Slide = slides.add_slide(layout)
        shapes: SlideShapes = new_slide.shapes
        placeholders: SlidePlaceholders = shapes.placeholders
        try:
            placeholders[0].text = line
        except Exception as e:
            raise RuntimeError(
                "Provided layout must use a textbox as the first placeholder"
            )
    return template
