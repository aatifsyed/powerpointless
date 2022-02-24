import logging
from typing import List

from pptx import Presentation as presentation
from pptx.presentation import Presentation
from pptx.shapes.shapetree import Shape
from pptx.slide import (
    Slide,
    SlideLayout,
    SlideMaster,
    SlidePlaceholders,
    Slides,
    SlideShapes,
)

logger = logging.getLogger(__name__)


def create_subtitles(template: Presentation, lines: List[str]) -> Presentation:
    try:
        slide_master: SlideMaster = template.slide_master
    except Exception as e:
        raise RuntimeError(
            "Provided presentation must have at least one slide master attached"
        ) from e

    try:
        layout: SlideLayout = slide_master.slide_layouts[0]
    except Exception as e:
        raise RuntimeError("Slide master must have at least one layout") from e

    slides: Slides = template.slides

    for line in lines:
        new_slide: Slide = slides.add_slide(layout)
        shapes: SlideShapes = new_slide.shapes
        placeholders: SlidePlaceholders = shapes.placeholders
        try:
            placeholders[0].text = line
        except Exception as e:
            raise RuntimeError(
                "Provided layout must use a textbox as the first placeholder"
            )
    return template


def extract_subtitles(source: Presentation) -> List[str]:
    lis: List[str] = []

    slides: Slides = source.slides
    slide: Slide
    for slide in slides:
        shapes: SlideShapes = slide.shapes
        shape: Shape
        for shape in shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                if len(text_frame.text) > 0:
                    lis.append(text_frame.text.rstrip())
    return lis
